import os
import docx2txt
import pandas as pd
from pptx import Presentation
from PyPDF2 import PdfFileReader
import pytesseract # will convert the image to text string
from PIL import Image # adds image processing capabilities


def getextension(f): return os.path.splitext(f)[1]


def read_txtfile(filename):
	"""To read ascii text file"""
	with open(filename) as file: lines = ''.join([line.rstrip() for line in file])
	return lines


def read_docxfile(filename):
	"""To read .docx file"""
	return [line for line in docx2txt.process(filename).splitlines() if line != '']


def read_xlsxfile(filename):
	"""To read .xlsx file"""
	file = pd.ExcelFile(filename)
	df = pd.concat([file.parse(name) for name in file.sheet_names], sort=False)
	return df.to_string(index=False)


def read_pptxfile(filename):
	"""To read .pptx file"""
	prs = Presentation(filename)
	return ' '.join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])


def read_pdffile(filename):
	"""To read .pdf file"""
	with open(filename, 'rb') as pdffile:
		pdfReader = PdfFileReader(pdffile)
		text = ' '.join(sum([pdfReader.getPage(n).extractText().splitlines() for n in range(pdfReader.numPages)], []))
	return text


def read_jpgfile(filename):
	"""To read .jpg file"""
	return pytesseract.image_to_string(Image.open(filename), lang='eng')


def read_file(filename):
	"""To read file as text"""
	reader = {
		  '.txt'  : read_txtfile
		, '.csv'  : read_txtfile
		, '.docx' : read_docxfile
		, '.xlsx' : read_xlsxfile
		, '.pptx' : read_pptxfile
		, '.jpg'  : read_jpgfile
		, '.pdf'  : read_pdffile
	}
	f_ext = getextension(filename)
	if f_ext not in reader.keys(): raise Exception("Not able to read " + f_ext)
	return reader[f_ext](filename)


	
	